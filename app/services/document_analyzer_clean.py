#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Service d'analyse de documents CV - Version améliorée
"""

import os
import re
from typing import Any
from typing import Dict
from typing import List

import pdfplumber
import pypdf as PyPDF2
import streamlit as st
from docx import Document
from pptx import Presentation


class DocumentAnalyzer:
    """Service d'analyse et d'extraction d'informations des CV"""

    # Clients connus pour améliorer la détection
    CLIENTS_CONNUS = [
        "BNP Paribas",
        "Société Générale",
        "Crédit Agricole",
        "BPCE",
        "Natixis",
        "AXA",
        "CNP Assurances",
        "Generali",
        "Allianz",
        "Swiss Life",
        "Orange",
        "SFR",
        "Bouygues Telecom",
        "Free",
        "Capgemini",
        "Accenture",
        "Sopra Steria",
        "Atos",
        "CGI",
        "IBM",
        "TCS",
        "Renault",
        "PSA",
        "Stellantis",
        "Michelin",
        "Valeo",
        "Total",
        "Engie",
        "EDF",
        "Veolia",
        "Suez",
        "SNCF",
        "Airbus",
        "Safran",
        "Thales",
        "Dassault",
        "Amazon",
        "Microsoft",
        "Google",
        "Meta",
        "Apple",
        "La Poste",
        "Pôle Emploi",
        "CAF",
        "CPAM",
        "FNAC",
        "Carrefour",
        "Leclerc",
        "Auchan",
        "Quanteam",
        "Rainbow Partners",
    ]

    @staticmethod
    def extract_text_from_file(file_path: str) -> str:
        """Extrait le texte d'un fichier avec lecture complète de toutes les pages"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()

            if file_extension == ".pdf":
                return DocumentAnalyzer._extract_text_from_pdf(file_path)
            elif file_extension in [".docx", ".doc"]:
                return DocumentAnalyzer._extract_text_from_docx(file_path)
            elif file_extension in [".pptx", ".ppt"]:
                return DocumentAnalyzer._extract_text_from_pptx(file_path)
            else:
                raise ValueError(f"Format de fichier non supporté: {file_extension}")

        except (OSError, ValueError, TypeError, AttributeError) as e:
            st.error(f"Erreur lors de l'extraction du fichier {file_path}: {e}")
            return ""

    @staticmethod
    def _extract_text_from_pdf(file_path: str) -> str:
        """Extrait le texte d'un PDF avec lecture page par page"""
        text_parts = []

        try:
            # Méthode 1 : pdfplumber (plus fiable pour les tableaux)
            st.info("📄 Extraction PDF avec pdfplumber...")
            with pdfplumber.open(file_path) as pdf:
                st.info(f"📄 PDF contient {len(pdf.pages)} page(s)")

                for page_num, page in enumerate(pdf.pages, 1):
                    st.info(f"📖 Lecture page {page_num}...")
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"--- PAGE {page_num} ---\n{page_text}")

                        # Extraire aussi les tableaux
                        tables = page.extract_tables()
                        for table in tables:
                            table_text = "\n".join(
                                [
                                    "\t".join([cell or "" for cell in row])
                                    for row in table
                                ]
                            )
                            if table_text.strip():
                                text_parts.append(
                                    f"--- TABLEAU PAGE {page_num} ---\n{table_text}"
                                )

                    except (OSError, ValueError, TypeError) as e:
                        st.warning(f"⚠️ Erreur page {page_num}: {e}")
                        continue

                if text_parts:
                    st.success(
                        f"✅ {len(text_parts)} sections extraites avec pdfplumber"
                    )
                    return "\n\n".join(text_parts)

            # Méthode 2 : PyPDF2 (fallback)
            st.info("📄 Fallback avec PyPDF2...")
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                st.info(f"📄 PDF contient {len(pdf_reader.pages)} page(s)")

                for page_num, page in enumerate(pdf_reader.pages, 1):
                    st.info(f"📖 Lecture page {page_num}...")
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(f"--- PAGE {page_num} ---\n{page_text}")
                    except (OSError, ValueError, TypeError) as e:
                        st.warning(f"⚠️ Erreur page {page_num}: {e}")
                        continue

                if text_parts:
                    st.success(f"✅ {len(text_parts)} sections extraites avec PyPDF2")
                    return "\n\n".join(text_parts)

            st.error("❌ Aucun texte extrait du PDF")
            return ""

        except (OSError, ValueError, TypeError, AttributeError) as e:
            st.error(f"❌ Erreur extraction PDF: {e}")
            return ""

    @staticmethod
    def _extract_text_from_docx(file_path: str) -> str:
        """Extrait le texte d'un document Word avec tous les éléments"""
        text_parts = []

        try:
            st.info("📄 Extraction Word avec python-docx...")
            doc = Document(file_path)

            # Extraire les paragraphes
            st.info(f"📄 Document contient {len(doc.paragraphs)} paragraphe(s)")
            for i, paragraph in enumerate(doc.paragraphs, 1):
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
                    if i % 50 == 0:  # Feedback tous les 50 paragraphes
                        st.info(f"📖 Traitement paragraphe {i}...")

            # Extraire les tableaux
            st.info(f"📄 Document contient {len(doc.tables)} tableau(x)")
            for table_num, table in enumerate(doc.tables, 1):
                st.info(f"📊 Extraction tableau {table_num}...")
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    table_text.append("\t".join(row_text))

                if table_text:
                    text_parts.append(
                        f"--- TABLEAU {table_num} ---\n" + "\n".join(table_text)
                    )

            result = "\n\n".join(text_parts)
            st.success(
                f"✅ {len(text_parts)} éléments extraits ({
                    len(result)} caractères)"
            )
            return result

        except (OSError, ValueError, TypeError, AttributeError) as e:
            st.error(f"❌ Erreur extraction Word: {e}")
            return ""

    @staticmethod
    def _extract_text_from_pptx(file_path: str) -> str:
        """Extrait le texte d'une présentation PowerPoint"""
        text_parts = []

        try:
            st.info("📄 Extraction PowerPoint avec python-pptx...")
            prs = Presentation(file_path)

            st.info(f"📄 Présentation contient {len(prs.slides)} slide(s)")

            for slide_num, slide in enumerate(prs.slides, 1):
                st.info(f"📖 Lecture slide {slide_num}...")
                slide_text = []

                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)

                        # Extraire aussi le texte des tableaux
                        if hasattr(shape, "table"):
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    row_text.append(cell.text.strip())
                                table_text.append("\t".join(row_text))
                            if table_text:
                                slide_text.append(
                                    "--- TABLEAU ---\n" + "\n".join(table_text)
                                )
                    except (OSError, ValueError, TypeError, AttributeError) as e:
                        st.warning(
                            f"⚠️ Erreur lors du traitement d'une forme PowerPoint: {e}"
                        )
                        continue

                if slide_text:
                    text_parts.append(
                        f"--- SLIDE {slide_num} ---\n" + "\n".join(slide_text)
                    )

            result = "\n\n".join(text_parts)
            st.success(
                f"✅ {len(text_parts)} slides traités ({len(result)} caractères)"
            )
            return result

        except (OSError, ValueError, TypeError, AttributeError) as e:
            st.error(f"❌ Erreur extraction PowerPoint: {e}")
            return ""

    @staticmethod
    def analyze_cv_content(text: str, consultant_name: str = "") -> Dict[str, Any]:
        """Analyse le contenu du CV et extrait les informations structurées"""
        st.info(f"🔍 Début de l'analyse pour {consultant_name}")

        analysis = {
            "consultant": consultant_name,
            "missions": [],
            "langages_techniques": [],
            "competences_fonctionnelles": [],
            "informations_generales": {},
            "texte_brut": text[:1000],  # Aperçu pour debug
        }

        try:
            # Extraction des missions
            missions = DocumentAnalyzer._extract_missions(text)
            analysis["missions"] = missions

            # Extraction des compétences techniques globales
            all_skills = set()
            for mission in missions:
                all_skills.update(mission.get("langages_techniques", []))

            # Ajouter les compétences globales du texte
            global_skills = DocumentAnalyzer._extract_technical_skills(text)
            all_skills.update(global_skills)

            analysis["langages_techniques"] = list(all_skills)[:25]  # Top 25

            # Extraction des compétences fonctionnelles
            analysis["competences_fonctionnelles"] = (
                DocumentAnalyzer._extract_functional_skills(text)
            )

            # Informations générales
            analysis["informations_generales"] = DocumentAnalyzer._extract_general_info(
                text
            )

            st.success(
                f"✅ Analyse terminée: {len(missions)} missions, {len(analysis['langages_techniques'])} technologies"
            )

        except (OSError, ValueError, TypeError, AttributeError, KeyError) as e:
            st.error(f"❌ Erreur analyse: {e}")
            import traceback

            traceback.print_exc()

        return analysis

    @staticmethod
    def _extract_missions(text: str) -> List[Dict]:
        """Extrait les missions du CV avec une approche multi-méthodes"""
        missions = []

        st.info(f"🔍 Analyse de {len(text)} caractères de texte...")

        # Nettoyer et préparer le texte
        text = text.replace("\r\n", "\n").replace("\r", "\n")

        # Séparer par pages pour mieux analyser
        pages = text.split("--- PAGE")
        st.info(f"📄 Analyse de {len(pages)} section(s) de document")

        all_text = " ".join(pages)  # Recombiner pour l'analyse globale

        # Trouver la section expérience
        experience_keywords = [
            "expérience professionnelle",
            "parcours professionnel",
            "historique professionnel",
            "expériences",
            "missions",
            "emplois",
            "postes",
            "carrière",
            "activités professionnelles",
            "réalisations",
            "projets",
            "interventions",
            "mandats",
            "consulting",
            "conseil",
        ]

        text_lower = all_text.lower()
        experience_start = 0

        for keyword in experience_keywords:
            match = re.search(keyword, text_lower)
            if match:
                experience_start = match.end()
                st.success(
                    f"✅ Section '{keyword}' trouvée à la position {experience_start}"
                )
                break

        # Prendre le texte d'expérience
        experience_text = (
            all_text[experience_start:] if experience_start > 0 else all_text
        )
        st.info(f"📝 Analyse de {len(experience_text)} caractères d'expérience")

        # Méthode 1: Analyse par blocs logiques
        missions.extend(DocumentAnalyzer._extract_missions_by_blocks(experience_text))

        # Méthode 2: Recherche par patterns spécifiques
        missions.extend(DocumentAnalyzer._extract_missions_by_patterns(experience_text))

        # Méthode 3: Recherche par clients connus
        missions.extend(
            DocumentAnalyzer._extract_missions_by_known_clients(experience_text)
        )

        st.info(f"🚀 {len(missions)} missions brutes trouvées avant nettoyage")

        # Nettoyer et dédupliquer
        unique_missions = DocumentAnalyzer._clean_and_deduplicate_missions(missions)

        # Trier par date de début (plus récent en premier)
        sorted_missions = sorted(
            unique_missions,
            key=lambda x: DocumentAnalyzer._date_sort_key(x.get("date_debut", "")),
            reverse=True,
        )

        st.success(f"✅ {len(sorted_missions)} missions finales après nettoyage")

        return sorted_missions[:20]  # Top 20

    @staticmethod
    def _extract_missions_by_blocks(text: str) -> List[Dict]:
        """Extraction par blocs de texte logiques"""
        missions = []

        # Diviser en blocs plus intelligemment
        blocks = re.split(r"\n\s*\n|\.\s*\n\s*\n|\.{2,}", text)

        # Filtrer les blocs significatifs (au moins 100 caractères)
        significant_blocks = [
            block.strip() for block in blocks if len(block.strip()) > 100
        ]

        st.info(f"📦 Analyse de {len(significant_blocks)} blocs significatifs")

        for i, block in enumerate(significant_blocks):
            mission = DocumentAnalyzer._extract_mission_from_block(
                block, block_num=i + 1
            )
            if mission and mission.get("client") and len(mission["client"]) > 2:
                missions.append(mission)

        return missions

    @staticmethod
    def _extract_mission_from_block(block: str, block_num: int = 0) -> Dict:
        """Extrait une mission d'un bloc de texte"""
        mission = {
            "date_debut": "",
            "date_fin": "",
            "client": "",
            "resume": "",
            "langages_techniques": [],
            "contexte": block[:500] + "..." if len(block) > 500 else block,
        }

        # Recherche de dates dans le bloc
        dates = DocumentAnalyzer._find_dates_in_text_improved(block)

        if len(dates) >= 2:
            mission["date_debut"] = dates[0]
            mission["date_fin"] = dates[1]
        elif len(dates) == 1:
            mission["date_debut"] = dates[0]
            mission["date_fin"] = "En cours"

        # Recherche de client
        client = DocumentAnalyzer._find_client_in_block_improved(block)
        if client:
            mission["client"] = client

        # Extraction du résumé long
        resume = DocumentAnalyzer._extract_long_mission_summary(block)
        mission["resume"] = resume

        # Extraction des technologies
        techs = DocumentAnalyzer._extract_technical_skills(block)
        mission["langages_techniques"] = techs[:10]

        return mission

    @staticmethod
    def _find_dates_in_text_improved(text: str) -> List[str]:
        """Trouve toutes les dates dans un texte - Version améliorée"""
        dates = []

        # Patterns de dates très complets
        date_patterns = [
            # Formats avec séparateurs
            (r"\b(\d{1,2})[\/\-\.](\d{1,2})[\/\-\.](\d{4})\b", "dmy"),  # DD/MM/YYYY
            (r"\b(\d{1,2})[\/\-\.](\d{4})\b", "my"),  # MM/YYYY
            (r"\b(\d{4})[\/\-\.](\d{1,2})\b", "ym"),  # YYYY/MM
            # Années dans un contexte temporel
            (r"\b(\d{4})\s*[-–—]\s*(\d{4})\b", "range"),  # 2020-2023
            # 2020-en cours
            (
                r"\b(\d{4})\s*[-–—]\s*(en\s+cours|actuel|présent|aujourd\'hui)\b",
                "ongoing",
            ),
            (r"\bdepuis\s+(\d{4})\b", "since"),  # depuis 2020
            (r"\ben\s+(\d{4})\b", "year"),  # en 2020
            # Mois en français
            (
                r"\b(janvier|février|mars|avril|mai|juin|juillet|août|septembre|octobre|novembre|décembre)\s+(\d{4})\b",
                "month_fr",
            ),
            (
                r"\b(jan|fév|mar|avr|mai|jun|jul|aoû|sep|oct|nov|déc)\.?\s+(\d{4})\b",
                "month_abbr",
            ),
        ]

        for pattern, pattern_type in date_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)

            for match in matches:
                if pattern_type == "dmy":
                    day, month, year = match.groups()
                    date_str = f"{year}-{month.zfill(2)}-{day.zfill(2)}"
                elif pattern_type == "my":
                    month, year = match.groups()
                    date_str = f"{year}-{month.zfill(2)}-01"
                elif pattern_type == "ym":
                    year, month = match.groups()
                    date_str = f"{year}-{month.zfill(2)}-01"
                elif pattern_type == "range":
                    start_year, end_year = match.groups()
                    dates.extend([f"{start_year}-01-01", f"{end_year}-12-31"])
                    continue
                elif pattern_type == "ongoing":
                    start_year = match.group(1)
                    dates.extend([f"{start_year}-01-01", "En cours"])
                    continue
                elif pattern_type in ["since", "year"]:
                    year = match.group(1)
                    date_str = f"{year}-01-01"
                elif pattern_type in ["month_fr", "month_abbr"]:
                    month_name, year = match.groups()
                    month_num = DocumentAnalyzer._month_name_to_number(
                        month_name.lower()
                    )
                    date_str = f"{year}-{month_num}-01"
                else:
                    continue

                if date_str and date_str not in dates:
                    dates.append(date_str)

        # Trier les dates et retourner
        valid_dates = [d for d in dates if d and d != "En cours"]
        valid_dates.sort(key=lambda x: DocumentAnalyzer._date_sort_key(x))

        # Ajouter "En cours" à la fin si présent
        if "En cours" in dates:
            valid_dates.append("En cours")

        return valid_dates[:6]  # Limiter à 6 dates max par bloc

    @staticmethod
    def _month_name_to_number(month_name: str) -> str:
        """Convertit un nom de mois français en numéro"""
        months = {
            "janvier": "01",
            "février": "02",
            "mars": "03",
            "avril": "04",
            "mai": "05",
            "juin": "06",
            "juillet": "07",
            "août": "08",
            "septembre": "09",
            "octobre": "10",
            "novembre": "11",
            "décembre": "12",
            "jan": "01",
            "fév": "02",
            "mar": "03",
            "avr": "04",
            "mai": "05",
            "jun": "06",
            "jul": "07",
            "aoû": "08",
            "sep": "09",
            "oct": "10",
            "nov": "11",
            "déc": "12",
        }
        return months.get(month_name, "01")

    @staticmethod
    def _find_client_in_block_improved(block: str) -> str:
        """Trouve le nom du client dans un bloc - Version améliorée"""
        # Patterns pour identifier les clients
        client_patterns = [
            r"(?:chez|pour|client[:\s]*)\s+([A-ZÀ-Ÿ][A-Za-zÀ-ÿ\s&\-\.]{3,60})",
            r"(?:société|entreprise|groupe|compagnie)\s+([A-ZÀ-Ÿ][A-Za-zÀ-ÿ\s&\-\.]{3,50})",
            r"([A-ZÀ-Ÿ][A-Za-zÀ-ÿ\s&\-\.]{3,50})\s+[-–—]\s*(?:consultant|développeur|chef|responsable|manager|directeur|lead|senior)",
            r"\b([A-Z][A-Za-z\s&\-\.]{3,50})\s+(?:SA|SAS|SARL|EURL|SNC|GIE|SCOP|AG|SE|SCA)\b",
            r"\b(Société\s+[A-ZÀ-Ÿ][A-Za-zÀ-ÿ\s&\-\.]{3,40})",
            r"\b(Groupe\s+[A-ZÀ-Ÿ][A-Za-zÀ-ÿ\s&\-\.]{3,40})",
            r"^([A-ZÀ-Ÿ][A-Za-zÀ-ÿ\s&\-\.]{5,50})\s*[-–—:]",
        ]

        for pattern in client_patterns:
            matches = re.finditer(pattern, block, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                client = match.group(1).strip()
                client = re.sub(r"\s*[-–—:]\s*.*$", "", client)
                client = re.sub(r"\s*\([^)]*\)\s*", "", client)
                client = DocumentAnalyzer._clean_client_name(client)

                if 3 < len(client) < 60:
                    return client

        # Chercher des clients connus
        for known_client in DocumentAnalyzer.CLIENTS_CONNUS:
            if known_client.lower() in block.lower():
                return known_client

        return ""

    @staticmethod
    def _extract_long_mission_summary(text: str) -> str:
        """Extrait un résumé long et détaillé de la mission (jusqu'à 1000 caractères)"""
        mission_keywords = [
            "mission",
            "projet",
            "développement",
            "conception",
            "réalisation",
            "mise en place",
            "création",
            "analyse",
            "étude",
            "design",
            "architecture",
            "implémentation",
            "déploiement",
            "maintenance",
            "support",
            "optimisation",
            "migration",
            "formation",
            "conseil",
            "audit",
            "expertise",
            "accompagnement",
            "pilotage",
        ]

        # Diviser en phrases
        sentences = re.split(r"[\.!?]+", text)
        relevant_sentences = []

        for sentence in sentences:
            sentence = sentence.strip()
            if 20 < len(sentence) < 500:
                sentence_lower = sentence.lower()
                score = 0

                # Score pour mots-clés
                for keyword in mission_keywords:
                    if keyword in sentence_lower:
                        score += 1

                # Score pour verbes d'action
                action_verbs = [
                    "développé",
                    "créé",
                    "conçu",
                    "réalisé",
                    "mis en place",
                    "déployé",
                ]
                for verb in action_verbs:
                    if verb in sentence_lower:
                        score += 3

                if score > 0:
                    relevant_sentences.append((sentence, score))

        # Construire le résumé
        if relevant_sentences:
            relevant_sentences.sort(key=lambda x: x[1], reverse=True)

            summary_parts = []
            total_length = 0
            max_length = 1000

            for sentence, score in relevant_sentences:
                clean_sentence = re.sub(r"^\s*[-–—•]\s*", "", sentence).strip()

                if total_length + len(clean_sentence) < max_length:
                    summary_parts.append(clean_sentence)
                    total_length += len(clean_sentence) + 2
                else:
                    break

            if summary_parts:
                full_summary = ". ".join(summary_parts)
                if not full_summary.endswith("."):
                    full_summary += "."
                return full_summary

        # Fallback
        clean_text = re.sub(r"^\s*[-–—•]\s*", "", text.strip())
        clean_text = re.sub(r"^\d{4}.*?:", "", clean_text)

        if len(clean_text) > 20:
            summary = clean_text[:1000]
            if len(clean_text) > 1000:
                last_period = summary.rfind(".")
                if last_period > 500:
                    summary = summary[: last_period + 1]
                else:
                    summary += "..."
            return summary

        return "Mission extraite automatiquement du CV"

    @staticmethod
    def _extract_missions_by_patterns(text: str) -> List[Dict]:
        """Recherche missions par patterns spécifiques"""
        missions = []

        # Pattern année-année avec description
        pattern1 = r"(\d{4})\s*[-–—]\s*(\d{4}|en\s+cours|actuel|présent)\s*[:\-]\s*([^\n\.]{30,})"
        matches1 = re.finditer(pattern1, text, re.IGNORECASE)

        for match in matches1:
            start_year, end_year, description = match.groups()
            mission = {
                "date_debut": f"{start_year}-01-01",
                "date_fin": (
                    "En cours"
                    if end_year.lower() in ["en cours", "actuel", "présent"]
                    else f"{end_year}-12-31"
                ),
                "client": DocumentAnalyzer._extract_client_from_text(description),
                "resume": description.strip()[:800]
                + ("..." if len(description) > 800 else ""),
                "langages_techniques": DocumentAnalyzer._extract_technical_skills(
                    description
                ),
            }
            if mission["client"]:
                missions.append(mission)

        return missions

    @staticmethod
    def _extract_missions_by_known_clients(text: str) -> List[Dict]:
        """Extraction basée sur les clients connus"""
        missions = []

        for client in DocumentAnalyzer.CLIENTS_CONNUS:
            for match in re.finditer(re.escape(client.lower()), text.lower()):
                start_pos = max(0, match.start() - 200)
                end_pos = min(len(text), match.end() + 300)
                context = text[start_pos:end_pos]

                dates = DocumentAnalyzer._find_dates_in_text_improved(context)
                if dates:
                    mission = {
                        "date_debut": dates[0] if dates else "",
                        "date_fin": dates[1] if len(dates) > 1 else "En cours",
                        "client": client,
                        "resume": DocumentAnalyzer._extract_long_mission_summary(
                            context
                        ),
                        "langages_techniques": DocumentAnalyzer._extract_technical_skills(
                            context
                        ),
                    }
                    missions.append(mission)

        return missions

    @staticmethod
    def _extract_client_from_text(text: str) -> str:
        """Extrait le nom du client d'un texte"""
        for client in DocumentAnalyzer.CLIENTS_CONNUS:
            if client.lower() in text.lower():
                return client

        patterns = [
            r"\b([A-Z][A-Za-z\s&]{3,40})\s+(?:SA|SAS|SARL|EURL|SNC|GIE)\b",
            r"(?:chez|pour|client)\s+([A-Z][A-Za-z\s&]{3,40})",
            r"\b([A-Z]{2,})\b",
        ]

        for pattern in patterns:
            match = re.search(pattern, text)
            if match:
                client = match.group(1).strip()
                if 3 < len(client) < 50:
                    return client

        return ""

    @staticmethod
    def _clean_and_deduplicate_missions(missions: List[Dict]) -> List[Dict]:
        """Nettoie et déduplique les missions"""
        unique_missions = []
        seen_combinations = set()

        for mission in missions:
            client = mission.get("client", "").strip()
            date_debut = mission.get("date_debut", "").strip()

            if not client or len(client) < 3:
                continue

            year = date_debut[:4] if len(date_debut) >= 4 else "unknown"
            key = f"{client.lower().strip()}_{year}"

            if key not in seen_combinations:
                seen_combinations.add(key)
                unique_missions.append(mission)

        return unique_missions

    @staticmethod
    def _date_sort_key(date_str: str) -> str:
        """Crée une clé de tri pour les dates"""
        if not date_str or date_str == "En cours":
            return "9999-12-31"

        if len(date_str) == 4:
            return f"{date_str}-01-01"
        elif len(date_str) >= 10:
            return date_str[:10]
        else:
            return date_str.ljust(10, "0")

    @staticmethod
    def _extract_technical_skills(text: str) -> List[str]:
        """Extrait les compétences techniques du texte"""
        skills = []
        text_lower = text.lower()

        # Langages de programmation
        programming_languages = [
            "java",
            "python",
            "javascript",
            "typescript",
            "c++",
            "c#",
            "php",
            "ruby",
            "go",
            "scala",
            "kotlin",
            "swift",
            "sql",
            "html",
            "css",
            "sass",
            "xml",
            "json",
        ]

        # Frameworks et technologies
        frameworks = [
            "react",
            "angular",
            "vue.js",
            "spring",
            "spring boot",
            "django",
            "flask",
            "express",
            "node.js",
            "laravel",
            "symfony",
            ".net",
            "asp.net",
        ]

        # Outils et plateformes
        tools = [
            "docker",
            "kubernetes",
            "jenkins",
            "gitlab",
            "github",
            "aws",
            "azure",
            "mongodb",
            "postgresql",
            "mysql",
            "oracle",
            "redis",
            "elasticsearch",
        ]

        all_skills = programming_languages + frameworks + tools

        for skill in all_skills:
            pattern = r"\b" + re.escape(skill.lower()) + r"\b"
            if re.search(pattern, text_lower):
                skills.append(skill.title())

        return list(set(skills))[:15]

    @staticmethod
    def _extract_functional_skills(text: str) -> List[str]:
        """Extrait les compétences fonctionnelles"""
        skills = []
        text_lower = text.lower()

        functional_skills = [
            "gestion de projet",
            "management",
            "leadership",
            "formation",
            "conseil",
            "analyse fonctionnelle",
            "architecture",
            "design",
            "scrum",
            "agile",
            "devops",
            "tests",
            "qualité",
            "sécurité",
            "performance",
        ]

        for skill in functional_skills:
            if skill in text_lower:
                skills.append(skill.title())

        return skills[:10]

    @staticmethod
    def _extract_general_info(text: str) -> Dict[str, str]:
        """Extrait les informations générales"""
        info = {}

        # Recherche email
        email_match = re.search(
            r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b", text
        )
        if email_match:
            info["email"] = email_match.group()

        # Recherche téléphone
        phone_match = re.search(r"(\+33|0)[1-9](?:[-.\s]?\d{2}){4}", text)
        if phone_match:
            info["telephone"] = phone_match.group()

        return info

    @staticmethod
    def _clean_client_name(client: str) -> str:
        """Nettoie le nom du client"""
        client = re.sub(r"[^\w\s&\-àâäéèêëïîôùûüÿç]", "", client)
        client = " ".join(client.split())

        if len(client) > 2:
            client = client.title()

        return client
