"""
Service de gestion des documents et CV
Gère l'upload, le stockage et l'analyse des documents
"""

import os
import shutil
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any
from typing import Dict
from typing import List
from typing import Optional

import streamlit as st

# Import pour le parsing de documents
try:
    import pdfplumber
    import PyPDF2
    from docx import Document as DocxDocument
    from pptx import Presentation
except ImportError as e:
    st.error(f"❌ Dépendance manquante pour le parsing de documents: {e}")

from database.database import get_database_session
from database.models import Consultant

class DocumentService:
    """Service pour la gestion des documents"""

    UPLOAD_DIR = Path("data/uploads")
    ALLOWED_EXTENSIONS = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'doc': 'application/msword',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'ppt': 'application/vnd.ms-powerpoint'}

    @classmethod
    def init_upload_directory(cls):
        """Initialise le répertoire d'upload"""
        cls.UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
        return cls.UPLOAD_DIR

    @classmethod
    def get_file_extension(cls, filename: str) -> str:
        """Récupère l'extension d'un fichier"""
        return filename.lower().split('.')[-1] if '.' in filename else ''

    @classmethod
    def is_allowed_file(cls, filename: str) -> bool:
        """Vérifie si le fichier est autorisé"""
        extension = cls.get_file_extension(filename)
        return extension in cls.ALLOWED_EXTENSIONS

    @classmethod
    def save_uploaded_file(cls, uploaded_file, consultant_id: int) -> Dict[str, Any]:
        """
        Sauvegarde un fichier uploadé

        Args:
            uploaded_file: Fichier Streamlit uploadé
            consultant_id: ID du consultant associé

        Returns:
            Dict avec les informations du fichier sauvegardé
        """
        try:
            # Initialiser le répertoire
            upload_dir = cls.init_upload_directory()

            # Créer un sous-répertoire pour le consultant
            consultant_dir = upload_dir / f"consultant_{consultant_id}"
            consultant_dir.mkdir(exist_ok=True)

            # Générer un nom de fichier unique
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            file_extension = cls.get_file_extension(uploaded_file.name)
            safe_filename = f"{timestamp}_{uploaded_file.name}"
            file_path = consultant_dir / safe_filename

            # Sauvegarder le fichier
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            # Retourner les informations
            return {
                'success': True,
                'file_path': str(file_path),
                'filename': safe_filename,
                'original_name': uploaded_file.name,
                'size': uploaded_file.size,
                'type': uploaded_file.type,
                'extension': file_extension,
                'consultant_id': consultant_id,
                'upload_date': datetime.now().isoformat()
            }

        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

    @classmethod
    def extract_text_from_file(cls, file_path: str) -> str:
        """Extrait le texte d'un fichier"""
        try:
            extension = cls.get_file_extension(file_path)

            if extension == 'pdf':
                return cls._extract_text_from_pdf(file_path)
            elif extension == 'docx':
                return cls._extract_text_from_docx(file_path)
            elif extension == 'pptx':
                return cls._extract_text_from_pptx(file_path)
            else:
                return f"Format {extension} non supporté"

        except Exception as e:
            return f"Erreur d'extraction: {e}"

    @classmethod
    def _extract_text_from_pdf(cls, file_path: str) -> str:
        """Extrait le texte d'un PDF"""
        try:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text.strip()
        except Exception as e:
            return f"Erreur PDF: {e}"

    @classmethod
    def _extract_text_from_docx(cls, file_path: str) -> str:
        """Extrait le texte d'un DOCX"""
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Erreur DOCX: {e}"

    @classmethod
    def _extract_text_from_pptx(cls, file_path: str) -> str:
        """Extrait le texte d'un PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Erreur PPTX: {e}"

    @classmethod
    def extract_text_from_pdf(cls, file_path: str) -> str:
        """Extrait le texte d'un fichier PDF"""
        try:
            text = ""

            # Méthode 1: pdfplumber (plus robuste)
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except BaseException:
                # Méthode 2: PyPDF2 (fallback)
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"

            return text.strip()

        except Exception as e:
            return f"Erreur lors de l'extraction PDF: {str(e)}"

    @classmethod
    def extract_text_from_docx(cls, file_path: str) -> str:
        """Extrait le texte d'un fichier Word DOCX"""
        try:
            doc = DocxDocument(file_path)
            text = ""

            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            # Extraire aussi le texte des tableaux
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"

            return text.strip()

        except Exception as e:
            return f"Erreur lors de l'extraction DOCX: {str(e)}"

    @classmethod
    def extract_text_from_pptx(cls, file_path: str) -> str:
        """Extrait le texte d'un fichier PowerPoint PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            return text.strip()

        except Exception as e:
            return f"Erreur lors de l'extraction PPTX: {str(e)}"

    @classmethod
    def extract_text_from_file(cls, file_path: str) -> str:
        """
        Extrait le texte d'un fichier selon son extension

        Args:
            file_path: Chemin vers le fichier

        Returns:
            Texte extrait du fichier
        """
        extension = cls.get_file_extension(file_path)

        if extension == 'pdf':
            return cls.extract_text_from_pdf(file_path)
        elif extension == 'docx':
            return cls.extract_text_from_docx(file_path)
        elif extension == 'pptx':
            return cls.extract_text_from_pptx(file_path)
        elif extension in ['doc', 'ppt']:
            return f"Format {
                extension.upper()} détecté mais nécessite une conversion vers {extension}x pour l'extraction de texte."
        else:
            return "Format de fichier non supporté pour l'extraction de texte."

    @classmethod
    def analyze_cv_content(cls, text: str) -> Dict[str, Any]:
        """
        Analyse le contenu d'un CV et extrait des informations

        Args:
            text: Texte extrait du CV

        Returns:
            Dict avec les informations analysées
        """
        analysis = {
            'skills_detected': [],
            'experience_years': {},
            'languages': [],
            'education': [],
            'contact_info': {},
            'summary': ""
        }

        text_lower = text.lower()

        # Détection de compétences techniques (liste non exhaustive)
        tech_skills = [
            'python',
            'java',
            'javascript',
            'typescript',
            'c++',
            'c#',
            'php',
            'ruby',
            'go',
            'rust',
            'react',
            'angular',
            'vue',
            'node.js',
            'express',
            'django',
            'flask',
            'spring',
            'sql',
            'mysql',
            'postgresql',
            'mongodb',
            'redis',
            'elasticsearch',
            'aws',
            'azure',
            'gcp',
            'docker',
            'kubernetes',
            'terraform',
            'git',
            'jenkins',
            'ci/cd',
            'devops',
            'machine learning',
            'deep learning',
            'ai',
            'data science',
            'analytics',
            'excel',
            'powerbi',
            'tableau',
            'looker',
            'agile',
            'scrum',
            'kanban',
            'project management']

        for skill in tech_skills:
            if skill.lower() in text_lower:
                analysis['skills_detected'].append({
                    'skill': skill.title(),
                    'confidence': 0.8,  # Confiance basique
                    'context_count': text_lower.count(skill.lower())
                })

        # Détection d'années d'expérience (patterns basiques)
        import re
        experience_patterns = [
            r'(\d+)\s*(?:ans?|années?)\s*(?:d\'?expérience|experience)',
            r'(\d+)\s*(?:years?)\s*(?:of\s*)?experience',
            r'expérience\s*(?:de\s*)?(\d+)\s*(?:ans?|années?)',
            r'experience\s*(?:of\s*)?(\d+)\s*years?'
        ]

        for pattern in experience_patterns:
            matches = re.findall(pattern, text_lower)
            for match in matches:
                years = int(match)
                if years <= 50:  # Validation basique
                    analysis['experience_years']['total'] = years
                    break

        # Détection de langues
        languages = [
            'français',
            'anglais',
            'espagnol',
            'allemand',
            'italien',
            'chinois',
            'japonais']
        for lang in languages:
            if lang in text_lower:
                analysis['languages'].append(lang.title())

        # Créer un résumé basique
        lines = text.split('\n')[:10]  # Premières lignes
        analysis['summary'] = ' '.join([line.strip()
                                       for line in lines if line.strip()])[:300] + "..."

        return analysis

    @classmethod
    def get_consultant_documents(cls, consultant_id: int) -> List[Dict[str, Any]]:
        """
        Récupère la liste des documents d'un consultant

        Args:
            consultant_id: ID du consultant

        Returns:
            Liste des documents
        """
        try:
            consultant_dir = cls.UPLOAD_DIR / f"consultant_{consultant_id}"
            documents = []

            if consultant_dir.exists():
                for file_path in consultant_dir.iterdir():
                    if file_path.is_file():
                        stat = file_path.stat()
                        documents.append({
                            'filename': file_path.name,
                            'path': str(file_path),
                            'size': stat.st_size,
                            'size_mb': round(stat.st_size / (1024 * 1024), 2),
                            'modified': datetime.fromtimestamp(stat.st_mtime),
                            'extension': cls.get_file_extension(file_path.name),
                            'type': cls.get_file_type_name(file_path.name)
                        })

            return sorted(documents, key=lambda x: x['modified'], reverse=True)

        except Exception as e:
            st.error(f"Erreur lors de la récupération des documents: {e}")
            return []

    @classmethod
    def get_file_type_name(cls, filename: str) -> str:
        """Retourne le nom du type de fichier"""
        extension = cls.get_file_extension(filename)
        type_names = {
            'pdf': 'PDF',
            'docx': 'Word (DOCX)',
            'doc': 'Word (DOC)',
            'pptx': 'PowerPoint (PPTX)',
            'ppt': 'PowerPoint (PPT)'
        }
        return type_names.get(extension, 'Inconnu')

    @classmethod
    def delete_document(cls, file_path: str) -> bool:
        """
        Supprime un document

        Args:
            file_path: Chemin vers le fichier

        Returns:
            True si suppression réussie
        """
        try:
            Path(file_path).unlink()
            return True
        except Exception as e:
            st.error(f"Erreur lors de la suppression: {e}")
            return False

    @classmethod
    def get_all_consultants_for_selection(cls) -> List[Dict[str, Any]]:
        """Récupère tous les consultants pour la sélection"""
        try:
            with get_database_session() as session:
                consultants = session.query(Consultant).all()
                return [
                    {
                        'id': consultant.id, 'name': f"{
                            consultant.prenom} {
                            consultant.nom}", 'email': consultant.email, 'display': f"{
                            consultant.prenom} {
                            consultant.nom} ({
                            consultant.email})"} for consultant in consultants]
        except Exception as e:
            st.error(f"Erreur lors de la récupération des consultants: {e}")
            return []