"""
Analyseur de CV simplifié et fonctionnel
"""
import re
from typing import Dict, List, Any
import streamlit as st

class SimpleDocumentAnalyzer:
    """Analyseur de CV simplifié"""
    
    # Technologies communes
    TECHNOLOGIES = [
        'Python', 'Java', 'JavaScript', 'TypeScript', 'C#', 'C++', 'PHP', 'Ruby', 'Go', 'Rust',
        'SQL', 'PostgreSQL', 'MySQL', 'Oracle', 'MongoDB', 'Redis', 'Elasticsearch',
        'React', 'Angular', 'Vue.js', 'Node.js', 'Express', 'Django', 'Flask', 'Spring',
        'Docker', 'Kubernetes', 'AWS', 'Azure', 'GCP', 'Jenkins', 'GitLab', 'Git',
        'Pandas', 'NumPy', 'TensorFlow', 'PyTorch', 'Scikit-learn', 'Jupyter',
        'Power BI', 'Tableau', 'Excel', 'Talend', 'Informatica', 'SSIS'
    ]
    
    # Clients communs
    CLIENTS = [
        'BNP Paribas', 'BNPP', 'Société Générale', 'SG', 'SGCIB', 
        'Crédit Agricole', 'BPCE', 'Natixis', 'AXA', 'Generali',
        'Orange', 'SFR', 'Bouygues', 'Free', 'Capgemini', 'Accenture',
        'Sopra Steria', 'Atos', 'CGI', 'IBM', 'Quanteam'
    ]
    
    @staticmethod
    def extract_text_from_file(file_path: str) -> str:
        """Extraction simple de texte depuis un fichier"""
        try:
            if file_path.lower().endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif file_path.lower().endswith('.pdf'):
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text()
                        return text
                except:
                    return "Erreur lors de l'extraction PDF"
            elif file_path.lower().endswith('.docx'):
                try:
                    from docx import Document
                    doc = Document(file_path)
                    text = []
                    for paragraph in doc.paragraphs:
                        text.append(paragraph.text)
                    return '\n'.join(text)
                except:
                    return "Erreur lors de l'extraction DOCX"
            elif file_path.lower().endswith(('.pptx', '.ppt')):
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return '\n'.join(text)
                except Exception as e:
                    return f"Erreur lors de l'extraction PowerPoint: {str(e)}"
            else:
                return "Format de fichier non supporté"
        except Exception as e:
            return f"Erreur d'extraction: {str(e)}"
    
    @staticmethod
    def analyze_cv_content(text: str, consultant_name: str = "") -> Dict[str, Any]:
        """Analyse simplifiée du contenu CV"""
        
        st.info(f"🔍 Analyse en cours pour {consultant_name}...")
        
        # Initialiser le résultat
        result = {
            "consultant": consultant_name,
            "missions": [],
            "langages_techniques": [],
            "competences_fonctionnelles": [],
            "informations_generales": {},
            "texte_brut": text[:500] + "..." if len(text) > 500 else text
        }
        
        try:
            # 1. Recherche des technologies
            technologies_found = []
            text_upper = text.upper()
            
            for tech in SimpleDocumentAnalyzer.TECHNOLOGIES:
                if tech.upper() in text_upper:
                    technologies_found.append(tech)
            
            result["langages_techniques"] = technologies_found[:15]  # Top 15
            
            # 2. Recherche des clients/missions
            missions = []
            for client in SimpleDocumentAnalyzer.CLIENTS:
                if client.upper() in text_upper:
                    # Créer une mission basique
                    mission = {
                        "client": client,
                        "titre": f"Mission chez {client}",
                        "description": f"Intervention chez {client}",
                        "langages_techniques": [tech for tech in technologies_found[:5]],
                        "duree": "Non spécifiée"
                    }
                    missions.append(mission)
            
            result["missions"] = missions[:10]  # Top 10 missions
            
            # 3. Compétences fonctionnelles basiques
            competences_func = []
            if "BI" in text_upper or "BUSINESS INTELLIGENCE" in text_upper:
                competences_func.append("Business Intelligence")
            if "DATA" in text_upper:
                competences_func.append("Analyse de données")
            if "PROJET" in text_upper or "MANAGEMENT" in text_upper:
                competences_func.append("Gestion de projet")
            if "CONSEIL" in text_upper or "CONSULTING" in text_upper:
                competences_func.append("Conseil")
            
            result["competences_fonctionnelles"] = competences_func
            
            # 4. Informations générales
            result["informations_generales"] = {
                "longueur_texte": len(text),
                "nombre_mots": len(text.split()),
                "technologies_detectees": len(technologies_found),
                "clients_detectes": len(missions)
            }
            
            st.success(f"✅ Analyse terminée: {len(missions)} missions, {len(technologies_found)} technologies")
            
        except Exception as e:
            st.error(f"❌ Erreur pendant l'analyse: {str(e)}")
            
        return result
